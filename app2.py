import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from io import BytesIO
import os
import json
import re
import base64
from pypdf import PdfReader
from pptx.enum.shapes import MSO_SHAPE
def main():
    # === AJOUTE ÇA APRÈS TES IMPORTS (juste après les imports) ===
    def checkbox_group_with_select_all(title, items, key_prefix):
        """Groupe de checkboxes avec bouton Tout cocher/décocher"""
        if not items:
            st.markdown(f"<small style='color:gray'>Aucun {title.lower()}</small>", unsafe_allow_html=True)
            return []

        # Clé pour le bouton "Tout"
        select_all_key = f"select_all_{key_prefix}"
        if select_all_key not in st.session_state:
            st.session_state[select_all_key] = True  # tout coché par défaut

        # Bouton Tout cocher/décocher
        col_title, col_btn = st.columns([3, 1])
        with col_title:
            st.markdown(f"**{title}**")
        with col_btn:
            if st.button("Tout", key=f"btn_{key_prefix}"):
                st.session_state[select_all_key] = not st.session_state[select_all_key]
                st.rerun()

        selected = []
        all_checked = st.session_state[select_all_key]

        for item in items:
            key = f"{key_prefix}_{hash(item) % 100000}"
            checked = st.checkbox(item, value=all_checked, key=key)
            if checked:
                selected.append(item)

        # Si tout est décoché manuellement → on désactive le "Tout"
        if not selected and all_checked:
            st.session_state[select_all_key] = False

        return selected

    def calculer_annees_experience(experiences_text):
        """
        GT Technologies 2025 – Calcul EXACT, sans arrondi bidon
        Retourne : année actuelle - première année trouvée dans les expériences
        """
        if not experiences_text or len(experiences_text.strip()) < 10:
            return 10  # valeur par défaut si rien

        import re
        from datetime import datetime

        # On extrait toutes les années 19xx ou 20xx
        annees = re.findall(r'\b(19\d{2}|20\d{2})\b', experiences_text)
        
        annees_valides = []
        for a in annees:
            try:
                year = int(a)
                if 1990 <= year <= datetime.now().year + 5:
                    annees_valides.append(year)
            except:
                pass

        if not annees_valides:
            return 10

        # On prend la PLUS ANCIENNE année
        premiere_annee = min(annees_valides)
        annees_reelles = datetime.now().year - premiere_annee

        # ON NE FAIT AUCUN ARRONDIS → on retourne le chiffre exact
        return annees_reelles

    # --- Configuration de la page Streamlit ---
    st.set_page_config(page_title="GT Technologies CV Builder",page_icon="image.png", layout="wide")

    if 'cv_data' not in st.session_state:
        st.session_state.cv_data = None

    # --- Logo ---
    logo_b64 = ""
    if os.path.exists("logo GTT nEw.png"):
        try:
            with open("logo GTT nEw.png", "rb") as f:
                logo_b64 = base64.b64encode(f.read()).decode()
        except Exception:
            pass

    # --- Style CSS ---
    st.markdown(f"""
    <style>
        .header {{background:linear-gradient(90deg,#5C2D91,#7D3F9D); padding:25px; text-align:center; border-radius:15px; margin-bottom:40px;}}
        .header img {{height:70px; background:white; padding:10px; border-radius:12px;}}
        .header h1 {{color:white; margin:15px 0 0; font-size:2.6rem; font-weight:900;}}
        .stButton>button {{background:#5C2D91 !important; color:white !important; font-weight:bold !important; border-radius:12px !important; padding:14px 32px !important; width: 100%;}}
        .block-container {{padding-top: 2rem;}}
        .checkbox-group {{border: 1px solid #ddd; padding: 10px; border-radius: 5px; margin-bottom: 10px;}}
        .checkbox-group h4 {{margin-top: 0;}}
        .slide-section {{margin-bottom: 20px;}}
    </style>
    """, unsafe_allow_html=True)

    if logo_b64:
        st.markdown(f"""
        <div class="header">
            <img src="data:image/png;base64,{logo_b64}">
            <h1>CV Builder – GT Technologies</h1>
        </div>
        """, unsafe_allow_html=True)
    else:
        st.markdown(f"""<div class="header"><h1>CV Builder – GT Technologies</h1></div>""", unsafe_allow_html=True)

    st.markdown("Génération du format Grant Thornton Technologies")
    st.divider()

    # --- Configuration de l'API Gemini ---
    api_key = st.secrets.get("GEMINI_API_KEY", None)
    MODEL = "gemini-2.5-flash"

    if api_key:
        genai.configure(api_key=api_key)

    # --- Données par défaut ---
    DEFAULT_DATA = {
        "NOM": "Prénom NOM",
        "Poste": "",
        "email": "",
        "DOMAINE_EXPERIENCE": "",
        "SECTEURS_EXPERIENCE": "",
        "PROFIL": "",
        "EXPERIENCES_PERTINENTES": "",
        "REFERENCES_PERTINENTES": "",
        "DIPLOMES_TEXTUELS": ""
    }
    # ============================================================
    # FONCTION ULTIME : LIT TOUS LES CV (PDF, DOCX, PPTX) MÊME AVEC TABLEAUX
    # ============================================================
    def extract_text(file):
        """
        Lit n'importe quel CV : PDF (texte/tableau/scanné), DOCX, PPTX
        Retourne toujours du texte propre, même sur les pires fichiers pourris
        """
        text = ""
        file_name = file.name.lower()
        file_bytes = file.read()
        file.seek(0)  # très important

        # ====================== DOCX ======================
        if file.type in ["application/vnd.openxmlformats-officedocument.wordprocessingml.document", "application/msword"] or file_name.endswith(".docx"):
            try:
                from docx import Document
                doc = Document(file)
                for para in doc.paragraphs:
                    if para.text.strip():
                        text += para.text + "\n"
                for table in doc.tables:
                    for row in table.rows:
                        row_text = " | ".join([cell.text.strip() for cell in row.cells if cell.text.strip()])
                        if row_text:
                            text += row_text + "\n"
                return text.strip()
            except Exception as e:
                st.warning(f"Erreur lecture DOCX : {e}")

        # ====================== PPTX ======================
        elif file.type in ["application/vnd.openxmlformats-officedocument.presentationml.presentation"] or file_name.endswith(".pptx"):
            try:
                from pptx import Presentation
                prs = Presentation(file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text") and shape.text.strip():
                            text += shape.text + "\n"
                        if hasattr(shape, "table"):
                            for row in shape.table.rows:
                                row_text = " | ".join([cell.text.strip() for cell in row.cells])
                                if row_text.strip():
                                    text += row_text + "\n"
                return text.strip() if text else "CV vide détecté dans le PPTX"
            except Exception as e:
                st.warning(f"Erreur lecture PPTX : {e}")

        # ====================== PDF – LA MÉGA BOMBE ======================
        else:  # PDF ou inconnu
            # Méthode 1 : PyMuPDF → le roi absolu (gère 95% des PDF même scannés avec OCR intégré)
            try:
                import fitz  # pymupdf
                doc = fitz.open(stream=file_bytes, filetype="pdf")
                for page in doc:
                    # Essai 1 : texte natif
                    page_text = page.get_text("text")
                    if len(page_text.strip()) > 50:
                        text += page_text + "\n"
                    else:
                        # Essai 2 : si texte faible → on active l'OCR (même sur image !)
                        pix = page.get_pixmap(dpi=300)
                        img_bytes = pix.tobytes("png")
                        import pytesseract
                        from PIL import Image
                        import io
                        ocr_text = pytesseract.image_to_string(Image.open(io.BytesIO(img_bytes)), lang='fra')
                        text += ocr_text + "\n"
                doc.close()
                if len(text) > 100:
                    return text.strip()
            except ImportError:
                st.error("Installe pymupdf et pytesseract pour lire TOUS les PDF : pip install pymupdf pytesseract pillow")
            except Exception as e:
                pass

            # Méthode 2 : pdfplumber → excellent sur les tableaux
            try:
                import pdfplumber
                with pdfplumber.open(file) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                        # Si pas de texte → on force les tableaux
                        for table in page.extract_tables():
                            for row in table:
                                clean_row = [cell.replace("\n", " ").strip() if cell else "" for cell in row]
                                text += " | ".join(clean_row) + "\n"
                if len(text) > 100:
                    return text.strip()
            except:
                pass

            # Méthode 3 : pypdf en dernier recours
            try:
                from pypdf import PdfReader
                reader = PdfReader(file)
                for page in reader.pages:
                    t = page.extract_text()
                    if t:
                        text += t + "\n"
                return text.strip()
            except:
                pass

        # Si tout a échoué
        return text.strip() if text else "AUCUN TEXTE EXTRAIT – Le fichier est probablement une image pure ou corrompu."


    def analyze_cv(cv_text):
        """Analyse le CV avec Gemini et retourne les données en JSON."""
        if not api_key:
            st.warning("Clé API Gemini non trouvée. Utilisation des données par défaut.")
            return DEFAULT_DATA.copy()

        prompt = f"""
        Tu es expert RH chez GT Technologies Sénégal.
        Analyse le CV et retourne STRICTEMENT un JSON avec les champs suivants :

        {{
        "NOM": "Prénom NOM",
        "Poste": "Titre du poste ou Dernier poste",
        "email": "email@example.com",
        "DOMAINE_EXPERIENCE": "Liste des domaines d'expertise, un par ligne, formatée avec • au début. Exemple: •Domaine1\\n•Domaine2",
        "SECTEURS_EXPERIENCE": "Liste des secteurs d'activité, un par ligne, formatée avec • au début. Exemple: •Secteur1\\n•Secteur2",
        "PROFIL": "Un résumé synthétique du profil en 4 à 6 phrases complètes. Utilise \\n pour séparer les phrases pour la lisibilité.",
        "EXPERIENCES_PERTINENTES": "Liste chronologique de TOUTES les expériences professionnelles, une par ligne, formatée avec • au début. Exemple: •Date : Poste - Entreprise\\n•Date : Poste - Entreprise",
        "REFERENCES_PERTINENTES": "Liste de TOUTES les références professionnelles ou organisations mentionnées dans le CV.
        - Si une mission ou un rôle associé est mentionné, ajoute-la après un double point.
        - Chaque ligne doit respecter le format exact : 'Référence : Mission(Année ou Période)'.
        - Extrais la période (Année, Années, ou Période Complète) si elle est associée à la référence et ajoute-la entre parenthèses à la fin.
        - Si aucune mission n'est mentionnée pour une référence, mets seulement le nom de la référence.
        Exemple: •CORAF : Audit du système d'information (2025)\\n•Banque de Dakar (BDK) : Audit de la plateforme technique (2024)",
        "DIPLOMES_TEXTUELS": "Liste de TOUS les diplômes, formations et certifications (CISA, ISO, etc.), une par ligne, formatée avec • au début. Exemple: •Master 2 - Université X\\n•Certification CISA"
        }}

        CV à analyser : {cv_text[:40000]}
        """

        try:
            model = genai.GenerativeModel(MODEL)
            resp = model.generate_content(prompt)
            raw = resp.text.strip().replace("```json", "").replace("```", "")

            json_match = re.search(r"\{.*\}", raw, re.DOTALL)
            if json_match:
                data = json.loads(json_match.group(0))
                for k, v in DEFAULT_DATA.items():
                    if k not in data or not data[k]:
                        data[k] = v
                return data
            else:
                raise ValueError("Aucun JSON valide trouvé.")
        except Exception as e:
            st.error(f"Erreur Gemini: {e}. Retour aux données par défaut.")
            return DEFAULT_DATA.copy()

    # --- Couleur mauve ---
    MAUVE = RGBColor(92, 45, 145)

    def format_block(tf, title, lines, font_size=10, alignment=None):
        """Remplit une zone de texte avec un titre et une liste à puces."""
        tf.clear()

        # Titre
        p_title = tf.add_paragraph()
        p_title.text = title
        p_title.font.bold = True
        p_title.font.size = Pt(14)
        p_title.font.color.rgb = MAUVE
        p_title.font.name = "Calibri"

        # Contenu
        for line in lines:
            if line.strip():
                p = tf.add_paragraph()
                p.text = line.strip().replace('•', '').strip()
                p.font.size = Pt(font_size)
                p.font.name = "Calibri"

                if alignment:
                    p.alignment = alignment

                if title != "PROFIL":
                    p.level = 0

    # --- Fonction pour afficher les checkboxes par élément ---
    def display_element_checkboxes(title, elements, session_key):
        """Affiche les éléments sous forme de cases à cocher."""
        st.markdown(f"""
        <div class="checkbox-group">
            <h4>{title}</h4>
        """, unsafe_allow_html=True)

        if session_key not in st.session_state:
            st.session_state[session_key] = {element: False for element in elements}

        selected_elements = []
        for element in elements:
            if st.checkbox(element, key=f"{session_key}_{element}", value=st.session_state[session_key].get(element, False)):
                selected_elements.append(element)

        st.session_state[session_key] = {element: element in selected_elements for element in elements}
        st.markdown("</div>", unsafe_allow_html=True)
        return selected_elements

    def fill_ppt_smart(prs, data, diploma_files, slide_elements, selected_domaines, selected_secteurs):
        MAUVE = RGBColor(92, 45, 145)

        # ===================================================================
        # 1. NOM + EMAIL + POSTE → partout, avec icônes
        # ===================================================================
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                for p in tf.paragraphs:
                    text = p.text
                    text_upper = text.upper()

                    if "NOM" in text_upper:
                        p.clear()
                        run = p.add_run()
                        run.text = f"NOM: {data['NOM']}"
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run.font.color.rgb = MAUVE
                        run.font.name = "Calibri"

                    if "@" in text or "email" in text.lower():
                        p.clear()
                        run = p.add_run()
                        run.text = data['email']
                        run.font.size = Pt(11)
                        run.font.color.rgb = MAUVE
                        run.font.name = "Calibri"

                    if "POSTE" in text_upper:
                        p.clear()
                        run = p.add_run()
                        run.text = f"Poste: {data['Poste']}"
                        run.font.bold = True
                        run.font.size = Pt(12)
                        run.font.color.rgb = MAUVE

        # ===================================================================
        # 2. PROFIL + EXPERIENCES + DIPLÔMES (par slide)
        # ===================================================================
        for slide_idx in range(3):
            if slide_idx >= len(prs.slides):
                continue
            slide = prs.slides[slide_idx]
            sel = slide_elements[slide_idx]

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                text_upper = tf.text.upper()

                # PROFIL
                if "PROFIL" in text_upper:
                    tf.clear()
                    p_title = tf.add_paragraph()
                    p_title.text = "PROFIL"
                    p_title.font.bold = True
                    p_title.font.size = Pt(14)
                    p_title.font.color.rgb = MAUVE
                    for line in data["PROFIL"].split('\n'):
                        if line.strip():
                            p = tf.add_paragraph()
                            p.text = line.strip()
                            p.font.size = Pt(10)
                            p.font.name = "Calibri"

                # EXPERIENCES PERTINENTES + DIPLÔMES
                if "EXPERIENCE" in text_upper and "PERTINENTE" in text_upper and len(tf.paragraphs) >= 5:
                    tf.clear()
                    p = tf.add_paragraph()
                    p.text = "EXPERIENCES PERTINENTES"
                    p.font.bold = True
                    p.font.size = Pt(14)
                    p.font.color.rgb = MAUVE
                    p.font.name = "Calibri"

                    for exp in sel["EXPERIENCES_PERTINENTES"]:
                        p = tf.add_paragraph()
                        p.text = exp.strip()
                        p.font.size = Pt(10)
                        p.font.name = "Calibri"

                    if sel["DIPLOMES_TEXTUELS"]:
                        p = tf.add_paragraph()
                        p.text = "Diplômes :"
                        p.font.bold = True
                        p.font.size = Pt(11)
                        p.font.color.rgb = MAUVE
                        for dip in sel["DIPLOMES_TEXTUELS"]:
                            p = tf.add_paragraph()
                            p.text = dip.strip()
                            p.font.size = Pt(10)
                            p.font.name = "Calibri"
                # RÉFÉRENCES PERTINENTES (dans la même boucle que les autres)
                if any(word in text_upper for word in ["RÉFÉRENCE", "REFERENCE", "RÉFÉRENCES", "REFERENCES"]):
                    tf.clear()

                    # Titre
                    p_title = tf.add_paragraph()
                    p_title.text = "REFERENCES PERTINENTES"
                    p_title.font.bold = True
                    p_title.font.size = Pt(14)
                    p_title.font.color.rgb = MAUVE
                    p_title.font.name = "Calibri"

                    # Les références cochées pour CETTE slide
                    for ref in sel["REFERENCES_PERTINENTES"]:
                        if ref.strip():
                            p = tf.add_paragraph()
                            p.text = ref.strip()
                            p.font.size = Pt(10)
                            p.font.name = "Calibri"
        # ===================================================================
        # 3. DOMAINES & SECTEURS → GLOBAUX, affichés sur TOUTES les slides
        # ===================================================================
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                tf = shape.text_frame
                text_upper = tf.text.upper()

                # DOMAINES D'EXPERIENCE
                if "DOMAINE" in text_upper and "EXPERIENCE" in text_upper:
                    tf.clear()
                    p_title = tf.add_paragraph()
                    p_title.text = "DOMAINE D'EXPERIENCE"
                    p_title.font.bold = True
                    p_title.font.color.rgb = MAUVE
                    for dom in selected_domaines:
                        p = tf.add_paragraph()
                        p.text = dom.strip()
                        p.font.size = Pt(9)
                        p.font.name = "Calibri"

                # SECTEURS D'EXPERIENCES
                if "SECTEUR" in text_upper and "EXPERIENCE" in text_upper:
                    tf.clear()
                    p_title = tf.add_paragraph()
                    p_title.text = "SECTEURS D'EXPERIENCES"
                    p_title.font.bold = True
                    p_title.font.color.rgb = MAUVE
                    for sec in selected_secteurs:
                        p = tf.add_paragraph()
                        p.text = sec.strip()
                        p.font.size = Pt(9)
                        p.font.name = "Calibri"

        # ===================================================================
        # 4. Images diplômes
        # ===================================================================
        if diploma_files:
            img_idx = 0
            for i in range(3, len(prs.slides)):
                slide = prs.slides[i]
                positions = [(Inches(0.6), Inches(6.4)), (Inches(3.8), Inches(6.4)),
                            (Inches(7.0), Inches(6.4)), (Inches(10.2), Inches(6.4))]
                for left, top in positions:
                    if img_idx >= len(diploma_files):
                        break
                    try:
                        slide.shapes.add_picture(BytesIO(diploma_files[img_idx].getvalue()),
                                                left, top, width=Inches(3.0), height=Inches(2.2))
                    except:
                        pass
                    img_idx += 1
    # --- Interface Streamlit ---
    TEMPLATE = "CV Papa Malick GUEYE Offre.pptx"

    if not os.path.exists(TEMPLATE):
        st.error("Template manquant → **CV Papa Malick GUEYE Offre.pptx** doit être dans le dossier !")
    else:
        st.success("Template chargé – prêt à l'emploi")

    c1, c2 = st.columns([2, 1])
    with c1:
        cv_file = st.file_uploader("CV source (PDF/DOCX)", type=["pdf", "docx"])
    with c2:
        diplomas = st.file_uploader("Diplômes (images/PDF)", type=["png", "jpg", "jpeg", "pdf"], accept_multiple_files=True)

    if cv_file and st.button("Analyser avec Gemini", type="primary"):
        with st.spinner("Analyse en cours..."):
            text = extract_text(cv_file)
            if text:
                result = analyze_cv(text)
                st.session_state.cv_data = result
                st.success("Analyse terminée avec succès !")
                st.balloons()
            else:
                st.error("Impossible d'extraire le texte du CV.")

    if st.session_state.cv_data:
        d = st.session_state.cv_data
        d = dict(d)

        st.subheader("🛠️ Vérification et Ajustement des Données")

        # Champs modifiables
        col1, col2, col3 = st.columns(3)
        with col1:
            d = dict(d)
            d["NOM"] = st.text_input("NOM", d["NOM"])
            d["Poste"] = st.text_input("Poste", d["Poste"])
            d["email"] = st.text_input("Email", d["email"])
        with col2:
            d["DOMAINE_EXPERIENCE"] = st.text_area("Domaines (• par ligne)", d["DOMAINE_EXPERIENCE"], height=160)
        with col3:
            d["SECTEURS_EXPERIENCE"] = st.text_area("Secteurs (• par ligne)", d["SECTEURS_EXPERIENCE"], height=160)

        # Profil modifiable
        st.markdown("""
        <div class="slide-section">
            <h4>Modification du Profil</h4>
        </div>
        """, unsafe_allow_html=True)
        d["PROFIL"] = st.text_area("PROFIL (modifiable)", d["PROFIL"], height=200)
    # === SÉLECTION GLOBALE DES DOMAINES ET SECTEURS (présents sur TOUTES les slides) ===
        st.markdown("### Sélection globale (présents sur toutes les slides)")

        col_dom, col_sec = st.columns(2)

        with col_dom:
            domaine_list = [l.strip() for l in d["DOMAINE_EXPERIENCE"].split('\n') if l.strip()]
            selected_domaines = checkbox_group_with_select_all("Domaines d'expérience", domaine_list, "global_domaine")

        with col_sec:
            secteur_list = [l.strip() for l in d["SECTEURS_EXPERIENCE"].split('\n') if l.strip()]
            selected_secteurs = checkbox_group_with_select_all("Secteurs d'expérience", secteur_list, "global_secteur")

        # Sélection des éléments par slide
        st.markdown("""
        <div class="slide-section">
            <h4>Sélection des éléments à afficher par slide</h4>
        </div>
        """, unsafe_allow_html=True)

        # Préparation des listes
        exp_list = [l for l in d["EXPERIENCES_PERTINENTES"].split('\n') if l.strip()]
        ref_list = [l for l in d["REFERENCES_PERTINENTES"].split('\n') if l.strip()]
        dip_list = [l for l in d["DIPLOMES_TEXTUELS"].split('\n') if l.strip()]

        slide_elements = []
        for slide_num in range(1, 4):
            st.markdown(f"### Slide {slide_num}")

            selected_experiences = checkbox_group_with_select_all(
                f"Expériences pour Slide {slide_num}", exp_list, f"exp_slide_{slide_num}"
            )

            selected_references = checkbox_group_with_select_all(
                f"Références pour Slide {slide_num}", ref_list, f"ref_slide_{slide_num}"
            )

            selected_diplomes = checkbox_group_with_select_all(
                f"Diplômes pour Slide {slide_num}", dip_list, f"dip_slide_{slide_num}"
            )

            slide_elements.append({
                "EXPERIENCES_PERTINENTES": selected_experiences,
                "REFERENCES_PERTINENTES": selected_references,
                "DIPLOMES_TEXTUELS": selected_diplomes
            })

        if st.button("Générer le CV GT Technologies", type="primary"):
            if not os.path.exists(TEMPLATE):
                st.error("Template manquant. Impossible de générer.")
            else:
                prs = Presentation(TEMPLATE)
                # ROND +XX ANS – VERSION ULTRA-SIMPLE & INDÉSTRUCTIBLE (aucune erreur possible)
                annees_final = calculer_annees_experience(d["EXPERIENCES_PERTINENTES"])
                texte_rond = f"+ {annees_final} ans\nd’expériences"

                for slide in prs.slides:

                    # position
                    left_violet = Inches(9.2)
                    top = Inches(0.2)
                    width = Inches(1.5)
                    height = Inches(1.5)

                    # 📌 1 — Cercle ORANGE derrière (créé en premier)
                    left_orange = left_violet - Inches(0.42)  # léger décalage à gauche

                    orange = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, left_orange, top, width, height
                    )
                    orange.fill.solid()
                    orange.fill.fore_color.rgb = RGBColor(242, 101, 34)     # orange
                    orange.line.color.rgb = RGBColor(242, 101, 34)

                    # 📌 2 — Cercle VIOLET devant
                    violet = slide.shapes.add_shape(
                        MSO_SHAPE.OVAL, left_violet, top, width, height
                    )
                    violet.fill.solid()
                    violet.fill.fore_color.rgb = RGBColor(92, 45, 145)      # violet
                    violet.line.color.rgb = RGBColor(92, 45, 145)

                    # 📌 texte sur le cercle violet
                    tf = violet.text_frame
                    tf.clear()
                    p = tf.paragraphs[0]
                    p.text = texte_rond
                    p.font.name = "Calibri"
                    p.font.size = Pt(14)
                    p.font.bold = True
                    p.font.color.rgb = RGBColor(255, 255, 255)
                    p.alignment = PP_ALIGN.CENTER

                    tf.margin_top = 0
                    tf.margin_bottom = 0
                
                st.success("Rond orange en haut / mauve en bas parfait → {texte_rond}")
                fill_ppt_smart(prs, d, diplomas or [], slide_elements, selected_domaines, selected_secteurs)

                # Sauvegarde
                out = BytesIO()
                prs.save(out)
                out.seek(0)

                st.success("CV généré avec succès !")

                file_name_clean = f"CV_{d['NOM'].replace(' ', '_')}_GT_Technologies.pptx"

                st.download_button(
                    "Télécharger le CV",
                    data=out,
                    file_name=file_name_clean,
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation"
                )
                st.balloons()
if __name__ == "__main__":
    pass  # laisse vide si tout ton code est déjà en haut
