import os
import io
import re
import json
import datetime
import base64
from typing import Dict, Any
import streamlit as st

# --- Modules Word (python-docx)
from docx import Document
from docx.oxml import OxmlElement
from docx.oxml.ns import qn
from docx.shared import RGBColor
from docx.shared import Pt 
from docx.oxml import parse_xml
from docx.oxml.ns import nsdecls

# --- Modules optionnels (Vérification des dépendances)
try:
    import pdfplumber
except ImportError:
    pdfplumber = None

try:
    # Nécessaire pour la génération PPTX
    from pptx import Presentation
    from pptx.util import Inches, Pt 
    from pptx.dml.color import RGBColor as PptxRGBColor # Couleur spécifique PPTX
except ImportError:
    Presentation = None
    PptxRGBColor = None 


# --- Remplacement d'OpenAI par Google GenAI
try:
    from google import genai
    from google.genai import types
    GEMINI_AVAILABLE = True
except ImportError:
    genai = None
    types = None
    GEMINI_AVAILABLE = False

import google.generativeai as genai
def main():
    # Configuration unique et définitive (identique aux autres apps)
    try:
        genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    except Exception as e:
        st.error("Clé API Gemini manquante ou invalide dans secrets.toml")
        st.stop()    

    # =====================================================
    # ========== CONSTANTES & OUTILS ======================
    # =====================================================

    # Constantes de Style
    COLOR_PRIMARY = "#6A3297" # Violet/Pourpre (Pour CSS Streamlit)
    if PptxRGBColor:
        COLOR_PRIMARY_PPTX = PptxRGBColor(0x6A, 0x32, 0x97)
    else:
        COLOR_PRIMARY_PPTX = None


    # --- OUTILS DOCX ---

    def set_cell_shading(cell, color_hex="E6E0ED"):
        """Applique un fond de couleur (ombrage) à une cellule de tableau."""
        tc = cell._tc
        tcPr = tc.get_or_add_tcPr()
        for shd in tcPr.xpath('w:shd'):
            tcPr.remove(shd)
        shd_xml = parse_xml(
            f'<w:shd {nsdecls("w")} w:val="clear" w:color="auto" w:fill="{color_hex}"/>'
        )
        tcPr.append(shd_xml)


    def shade_header_row(table, color_hex="E6E0ED"):
        """Applique la couleur mauve clair à la première ligne d'un tableau."""
        if not table.rows:
            return
        header_row = table.rows[0]
        for cell in header_row.cells:
            set_cell_shading(cell, color_hex)
            for paragraph in cell.paragraphs:
                for run in paragraph.runs:
                    run.bold = True
                    run.font.name = "Trebuchet MS" 
                    run.font.size = Pt(11)
                    run.font.color.rgb = RGBColor(0, 0, 0)
                    run.font.name = "Trebuchet MS"
                    run._element.rPr.rFonts.set(qn('w:eastAsia'), 'Trebuchet MS')
    def set_run_black_bold(run):
        run.bold = True
        run.font.color.rgb = RGBColor(0, 0, 0)
        run.font.name = "Trebuchet MS" 
        run.font.size = Pt(11)
        # Force la police Trebuchet MS correctement
        run.font.name = "Trebuchet MS"
        r = run._element
        r.rPr.rFonts.set(qn('w:eastAsia'), 'Trebuchet MS')



    # =====================================================
    # ========== EXTRACTION DU TEXTE ======================
    # =====================================================

    def extract_from_docx_bytes(b: bytes) -> str:
        """Extrait le texte et les tables d'un fichier DOCX en bytes."""
        doc = Document(io.BytesIO(b))
        parts = []
        for p in doc.paragraphs:
            if p.text.strip():
                parts.append(p.text)
        for t in doc.tables:
            for r in t.rows:
                row_text = " | ".join(c.text for c in r.cells)
                if row_text.strip():
                    parts.append(row_text)
        return "\n".join(parts)

    def extract_from_pdf_bytes(b: bytes) -> str:
        """Extrait le texte d'un fichier PDF en bytes (nécessite pdfplumber)."""
        if not pdfplumber:
            raise RuntimeError("Installe pdfplumber : pip install pdfplumber")
        parts = []
        with pdfplumber.open(io.BytesIO(b)) as pdf:
            for page in pdf.pages:
                parts.append(page.extract_text() or "")
        return "\n".join(parts)

    def extract_from_pptx_bytes(b: bytes) -> str:
        """Extrait le texte d'un fichier PPTX en bytes (nécessite python-pptx)."""
        if not Presentation:
            raise RuntimeError("Installe python-pptx : pip install python-pptx")
        prs = Presentation(io.BytesIO(b))
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text)
        return "\n".join(texts)


    # =====================================================
    # ========== PROMPT ET APPEL À L’IA (GEMINI) ==========
    # =====================================================

    def build_prompt(text: str) -> str:
        """Construit le prompt pour Gemini avec consignes précises pour tout extraire."""
        short = text[:40000] if len(text) > 40000 else text

        return f"""
    Tu es un expert RH chargé d'extraire les informations détaillées d'un CV.

    ⚠️ Instructions importantes :
    - Pour la section 'experience' (expérience professionnelle), extrais uniquement les champs 'date', 'entreprise' et 'poste'.
    - Pour la section 'taches_detaillees' et les 'activites' dans les 'projets', **tu dois inclure toutes les tâches et activités présentes dans le CV**, même si elles semblent redondantes ou longues.
    - Ne résume pas et ne reformule pas les phrases : conserve la formulation originale des tâches dans la mesure du possible.
    - Si le CV contient plusieurs projets ou expériences détaillées, liste-les toutes dans la section 'projets' avec leurs activités respectives.
    - Le 'nom_candidat' correspond à la firme ou société qui présente le CV.
    - Le 'nom_employe' correspond au nom de la personne physique.

    Réponds UNIQUEMENT avec un JSON valide, sans texte explicatif, dans cette structure exacte :
    {{
        "poste": "",
        "nom_candidat": "",
        "nom_employe": "",
        "naissance": "",
        "nationalite": "",
        "pays_travailles": ["", ""],
        "education": [
            {{"date": "", "detail": "", "lieu": ""}}
        ],
        "autres_formations": [
            {{"date": "", "detail": "", "lieu": ""}}
        ],
        "langues": [
            {{"langue": "", "oral": "", "lu": "", "ecrit": ""}}
        ],
        "experience": [
            {{"date": "", "entreprise": "", "poste": ""}}
        ],
        "taches_detaillees": ["", ""],
        "projets": [
            {{
                "nom": "",
                "annee": "",
                "lieu": "",
                "poste": "",
                "caracteristiques": "",
                "activites": ["", ""]
            }}
        ]
    }}

    Voici le texte du CV :
    {short}
    """

    def parse_with_llm(text: str, api_key: str) -> Dict[str, Any]:
        """Appelle Gemini pour extraire les données en format JSON selon le schéma corrigé."""
        if not GEMINI_AVAILABLE:
            return {"_error": "Module google-genai non installé. Installez-le avec : pip install google-genai"}
        if not api_key:
            return {"_error": "Clé API Gemini manquante. Veuillez la configurer dans .streamlit/secrets.toml"}
        
        # Définition du schéma JSON (complet)
        schema = types.Schema(
            type=types.Type.OBJECT,
            properties={
                "poste": types.Schema(type=types.Type.STRING),
                "nom_candidat": types.Schema(type=types.Type.STRING),
                "nom_employe": types.Schema(type=types.Type.STRING),
                "naissance": types.Schema(type=types.Type.STRING),
                "nationalite": types.Schema(type=types.Type.STRING),
                "pays_travailles": types.Schema(
                    type=types.Type.ARRAY, items=types.Schema(type=types.Type.STRING)
                ),
                "education": types.Schema(
                    type=types.Type.ARRAY,
                    items=types.Schema(
                        type=types.Type.OBJECT,
                        properties={"date": types.Schema(type=types.Type.STRING), "detail": types.Schema(type=types.Type.STRING), "lieu": types.Schema(type=types.Type.STRING)},
                    ),
                ),
                "autres_formations": types.Schema(
                    type=types.Type.ARRAY,
                    items=types.Schema(
                        type=types.Type.OBJECT,
                        properties={"date": types.Schema(type=types.Type.STRING), "detail": types.Schema(type=types.Type.STRING), "lieu": types.Schema(type=types.Type.STRING)},
                    ),
                ),
                "langues": types.Schema(
                    type=types.Type.ARRAY,
                    items=types.Schema(
                        type=types.Type.OBJECT,
                        properties={"langue": types.Schema(type=types.Type.STRING), "oral": types.Schema(type=types.Type.STRING), "lu": types.Schema(type=types.Type.STRING), "ecrit": types.Schema(type=types.Type.STRING)},
                    ),
                ),
                "experience": types.Schema(
                    type=types.Type.ARRAY,
                    items=types.Schema(
                        type=types.Type.OBJECT,
                        properties={"date": types.Schema(type=types.Type.STRING), "entreprise": types.Schema(type=types.Type.STRING), "poste": types.Schema(type=types.Type.STRING)},
                    ),
                ),
                "taches_detaillees": types.Schema(
                    type=types.Type.ARRAY, items=types.Schema(type=types.Type.STRING)
                ),
                "projets": types.Schema(
                    type=types.Type.ARRAY,
                    items=types.Schema(
                        type=types.Type.OBJECT,
                        properties={"nom": types.Schema(type=types.Type.STRING), "annee": types.Schema(type=types.Type.STRING), "lieu": types.Schema(type=types.Type.STRING), "poste": types.Schema(type=types.Type.STRING), "caracteristiques": types.Schema(type=types.Type.STRING), "activites": types.Schema(type=types.Type.ARRAY, items=types.Schema(type=types.Type.STRING))},
                    ),
                ),
            },
            required=["poste", "nom_candidat", "nom_employe", "education", "experience", "projets"],
        )

        # Configuration unique (une seule fois dans tout le fichier, juste après les imports)
    try:
        genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    except:
        pass

    GEMINI_MODEL = "gemini-2.5-flash"   # ou "gemini-1.5-flash" si tu veux encore plus de quota gratuit


    def parse_with_llm(text: str) -> dict:
        """Version ultra-légère, identique à tes autres apps qui marchent parfaitement"""
        
        if not text.strip():
            return {"_error": "Aucun texte extrait"}

        prompt = build_prompt(text[:40000])   # garde ta fonction build_prompt existante

        try:
            model = genai.GenerativeModel(
                model_name=GEMINI_MODEL,
                generation_config={
                    "temperature": 0.0,
                    "response_mime_type": "application/json"
                }
            )

            response = model.generate_content(prompt)
            raw = response.text.strip()

            # Nettoie les ```json que Gemini ajoute parfois Gemini
            if raw.startswith("```json"):
                raw = raw[7:]
            if raw.endswith("```"):
                raw = raw[:-3]

            return json.loads(raw)

        except Exception as e:
            msg = str(e).lower()
            if "quota" in msg or "429" in msg:
                return {"_error": "Quota Gemini dépassé (attends ou passe en payant)"}
            elif "10053" in msg:
                return {"_error": "Connexion coupée par ton antivirus/firewall → désactive Windows Defender 2 min"}
            else:
                return {"_error": f"Erreur Gemini : {e}"}

    # =====================================================
    # ========== GÉNÉRATION DU DOCX STANDARD =============
    # =====================================================

    def build_standard_docx(data: Dict[str, Any], buffer: io.BytesIO):
        doc = Document()
        today = datetime.datetime.now().strftime("%d/%m/%Y")
        section_counter = 1

        # 1. Poste
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Poste : ")
        set_run_black_bold(run)
        p.add_run(data.get('poste',''))
        section_counter += 1

        # 2. Nom du Candidat
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Nom du Candidat : ")
        set_run_black_bold(run)
        p.add_run(data.get('nom_candidat',''))
        section_counter += 1

        # 3. Nom de l'Employé
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Nom de l'Employé : ")
        set_run_black_bold(run)
        p.add_run(data.get('nom_employe',''))
        section_counter += 1

        # 4. Date de naissance et Nationalité
        # 4. Date de naissance et Nationalité
        p = doc.add_paragraph()

        # Date de naissance en gras
        run = p.add_run(f"{section_counter}. Date de naissance : ")
        set_run_black_bold(run)

        # Valeur date de naissance normale
        p.add_run(data.get('naissance',''))

        # Ajouter un espace avant Nationalité
        p.add_run("        ")  # ou "\t" pour tabulation

        # Nationalité en gras
        run_nat = p.add_run("Nationalité : ")
        set_run_black_bold(run_nat)

        # Valeur nationalité normale
        p.add_run(data.get('nationalite',''))

        section_counter += 1

        # 5. Éducation
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Éducation")
        set_run_black_bold(run)
        section_counter += 1

        t = doc.add_table(rows=1, cols=2)
        t.style = "Table Grid"
        t.rows[0].cells[0].text = "Date"
        t.rows[0].cells[1].text = "Écoles et diplômes obtenus"
        shade_header_row(t) 
        for e in data.get("education", []):
            r = t.add_row().cells
            r[0].text = e.get("date", "")
            r[1].text = f"{e.get('detail','')} ({e.get('lieu','')})".strip(" ()")

        # 6. Autres Formations
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Autres Formations")
        set_run_black_bold(run)
        section_counter += 1

        t_form = doc.add_table(rows=1, cols=2)
        t_form.style = "Table Grid"
        t_form.rows[0].cells[0].text = "Date"
        t_form.rows[0].cells[1].text = "Institution et certifications obtenues"
        shade_header_row(t_form)
        for f in data.get("autres_formations", []):
            r = t_form.add_row().cells
            r[0].text = f.get("date","")
            r[1].text = f"{f.get('detail','')} ({f.get('lieu','')})".strip(" ()")

        # 7. Pays travaillés
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Pays où l’employé a travaillé")
        set_run_black_bold(run)
        section_counter += 1
        doc.add_paragraph(", ".join(data.get("pays_travailles", [])) if data.get("pays_travailles") else "Non spécifié.")

        # 8. Langues
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Langues")
        set_run_black_bold(run)
        section_counter += 1

        langs = data.get("langues", [])
        if langs:
            t3 = doc.add_table(rows=1, cols=4)
            t3.style = "Table Grid"
            hdr = ["Langue", "Parlée", "Lue", "Écrite"]
            for i, h in enumerate(hdr):
                t3.rows[0].cells[i].text = h
            shade_header_row(t3)
            for l in langs:
                r = t3.add_row().cells
                r[0].text = l.get("langue","")
                r[1].text = l.get("oral","")
                r[2].text = l.get("lu","")
                r[3].text = l.get("ecrit","")
        else:
            doc.add_paragraph("Non spécifié.")

        # 9. Expérience professionnelle
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Expérience professionnelle")
        set_run_black_bold(run)
        section_counter += 1

        exps = data.get("experience", [])
        t4 = doc.add_table(rows=1, cols=3)
        t4.style = "Table Grid"
        hdr = ["Dates de l'emploi", "Nom de la société", "Poste occupé"]
        for i, h in enumerate(hdr):
            t4.rows[0].cells[i].text = h
        shade_header_row(t4)
        for e in exps:
            r = t4.add_row().cells
            r[0].text = e.get("date","")
            r[1].text = e.get("entreprise","")
            r[2].text = e.get("poste","")

        # Ajouter un paragraphe vide pour espacer le tableau
        doc.add_paragraph("")
        doc.add_paragraph("")

        # Tableau fusionnant les deux concepts
        all_projects = data.get("projets", [])
        t5 = doc.add_table(rows=1, cols=2)
        t5.style = "Table Grid"

        # --- En-têtes personnalisés
        r_hdr = t5.rows[0].cells
        r_hdr[0].text = "10. Détail des tâches exécutées"
        r_hdr[1].text = "11. Illustration de l'expérience"
        shade_header_row(t5)

        # --- Remplissage du tableau
        for i, p in enumerate(all_projects):
            r = t5.add_row().cells
            acts = p.get("activites", [])

            # Colonne gauche : activités (tâches)
            if acts:
                r[0].paragraphs[0].clear()
                for a in acts:
                    p_act = r[0].add_paragraph(style='List Bullet')
                    p_act.add_run(a)
            else:
                r[0].text = "XXXXXXXXXX"
                

            # Colonne droite : détails du projet
            details_col = r[1].paragraphs[0]
            details_col.text = ""
            # Texte avec seulement le label en gras
            details_col.add_run("Nom du projet ou de la mission : ").bold = True
            details_col.add_run(f"{p.get('nom','')}\n")
            details_col.add_run("Année : ").bold = True
            details_col.add_run(f"{p.get('annee','')}\n")

            details_col.add_run("Lieu : ").bold = True
            details_col.add_run(f"{p.get('lieu','')}\n")

            details_col.add_run("Principales caractéristiques : ").bold = True
            details_col.add_run(f"{p.get('caracteristiques','')}\n")

            details_col.add_run("Poste : ").bold = True
            details_col.add_run(f"{p.get('poste','')}")


        # 12. Attestation
        p = doc.add_paragraph()
        run = p.add_run(f"{section_counter}. Attestation")
        set_run_black_bold(run)
        section_counter += 1

        # Texte explicatif avant le tableau
        doc.add_paragraph(
            "Je soussigné certifie que les renseignements ci-dessus rendent fidèlement "
            "compte de ma situation, de mes qualifications et de mon expérience."
        )

        # Tableau d’attestation
        t_attest = doc.add_table(rows=2, cols=3)
        t_attest.style = "Table Grid"
        hdr_attest = t_attest.rows[0].cells
        hdr_attest[0].text = "Nom du représentant habilité"
        hdr_attest[1].text = "Signature"
        hdr_attest[2].text = "Date"
        shade_header_row(t_attest)

        data_attest = t_attest.rows[1].cells
        data_attest[0].text = data.get("nom_candidat", "Inconnu")
        data_attest[1].text = ""  # signature vide
        data_attest[2].text = today


        doc.save(buffer)
        buffer.seek(0)



    # =====================================================
    # ========== INTERFACE STREAMLIT (LE BAS DU SCRIPT) ===
    # =====================================================

    st.set_page_config(page_title="CV → Format GT", page_icon="image.png", layout="wide")

    def get_base64_of_image(image_file):
        if not os.path.exists(image_file):
            # Utiliser un logo par défaut ou un espace réservé si le fichier est manquant
            return "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII=" # Placeholder
        with open(image_file, "rb") as f: 
            data = f.read()
        return base64.b64encode(data).decode()

    LOGO_PATH = "logo GTT nEw.png" 
    COLOR_PRIMARY = "#6A3297" 
    COLOR_SECONDARY = "#828D4D" 
    COLOR_BACKGROUND_ACCENT = "#F3E6FF" 
    COLOR_BACKGROUND = "#f7f7f7" 

    try:
        # Utilisez get_base64_of_image(LOGO_PATH) dans un environnement réel.
        # Ici on utilise un placeholder par défaut pour assurer l'exécution.
        logo_base64 = get_base64_of_image(LOGO_PATH)
    except Exception as e:
        logo_base64 = "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAYAAAAfFcSJAAAADUlEQVR42mNkYAAAAAYAAjCB0C8AAAAASUVORK5CYII="


    # --- CSS GLOBALE ET DE MARQUE ---
    if logo_base64:
        st.markdown(f"""
        <style>
        /* 1. Reset Streamlit Defaults */
        #MainMenu {{visibility: hidden;}}
        footer {{visibility: hidden;}}
        header {{visibility: hidden;}}

        /* 2. Style de fond général */
        .stApp {{
            background-color: {COLOR_BACKGROUND};
            color: #333;
            font-family: 'Segoe UI', sans-serif;
            padding-top: 30px; 
        }}

        /* 3. Logo en haut à gauche (Position Fixe) */
        .logo-fixed-top-left {{
            position: fixed;
            top: 10px;
            left: 20px;
            z-index: 9999;
            background-color: {COLOR_BACKGROUND};
            padding-right: 20px;
            padding-bottom: 5px;
            border-radius: 0 0 5px 0;
        }}
        .logo-fixed-top-left img {{
            max-height: 40px; 
            width: auto;
        }}

        /* 4. Titres (Couleur Primaire GTT) - CENTRÉS */
        h1 {{
            color: {COLOR_PRIMARY}; 
            font-weight: 800; 
            text-align: center; 
            margin-top: 0; 
        }}
        h2 {{
            color: #555555; 
            font-weight: 400;
            font-size: 1.5em;
            text-align: center; 
            margin-bottom: 2rem;
        }}
        h3 {{
            color: {COLOR_PRIMARY};
            font-weight: 600;
            border-bottom: 2px solid {COLOR_PRIMARY}20;
            padding-bottom: 0.2em;
            text-align: center; 
            margin-top: 2rem;
        }}

        /* 5. Styles d'Upload - Titre visible et stylisé */
        div[data-testid="stFileUploader"] > div:first-child {{
            font-size: 1.2em;
            font-weight: 600;
            color: {COLOR_PRIMARY};
            text-align: center;
        }}

        div[data-testid="stFileUploader"] {{
            background-color: {COLOR_BACKGROUND_ACCENT};
            padding: 1.5em;
            border-radius: 12px;
            border: 2px dashed {COLOR_PRIMARY}40;
            margin-bottom: 1.5rem;
        }}

        /* --- TRADUCTION DE L'UPLOADEUR (CSS) --- */
        .stFileUploader > label + div > div:nth-child(2) > div:nth-child(1) > div:nth-child(1) > div:nth-child(1) > div:nth-child(1) > div:nth-child(1) > div:nth-child(1) > span:nth-child(1):empty::before {{
            content: "Glisser-déposer le fichier ici";
        }}
        .stFileUploader > label + div > div:nth-child(2) button span:empty::before {{
            content: "Parcourir les fichiers";
        }}
        .stFileUploader > label + div > div:nth-child(2) > div:nth-child(1) > div:nth-child(1) > div:nth-child(1) > div:nth-child(1) > div:nth-child(1) > div:nth-child(1) > span:nth-child(3):empty::before {{
            content: "ou";
        }}
        
        /* 6. Boutons et autres styles conservés */
        .stButton>button {{
            background-color: {COLOR_PRIMARY};
            color: white;
            border-radius: 8px;
            border: none;
            font-weight: 600;
            padding: 0.7em 1.5em;
            transition: background-color 0.3s;
        }}
        .stButton>button:hover {{
            background-color: #55277A; 
            box-shadow: 0 4px 8px rgba(0, 0, 0, 0.1);
        }}

        .stDownloadButton>button {{
            background-color: {COLOR_SECONDARY};
            color: white;
            font-weight: 600;
            border-radius: 8px;
            border: none;
        }}
        .stDownloadButton>button:hover {{
            background-color: #6C763F; 
        }}

        </style>
        """, unsafe_allow_html=True)
        
        # Affichage du Logo en haut à gauche
        st.markdown(f"""
            <div class="logo-fixed-top-left">
                <img src="data:image/png;base64,{logo_base64}" alt="Logo Grant Thornton Technologies">
            </div>
        """, unsafe_allow_html=True)


    if "results" not in st.session_state:
        st.session_state["results"] = []
    if "extracted_texts" not in st.session_state:
        st.session_state["extracted_texts"] = {}

    # --- TITRES PRINCIPAUX CENTRÉS ---
    st.markdown("<h1>Analyse & Standardisation CV</h1>", unsafe_allow_html=True)
    st.markdown("<h2>Génération du format Grant Thornton Technologies</h2>", unsafe_allow_html=True)

    # === LECTURE SÉCURISÉE DE LA CLÉ API ===
    #api_key = st.secrets.get("GEMINI_API_KEY") 

    #if not api_key:
        #st.error("⚠️ **Clé API Gemini Manquante.** Veuillez créer le fichier `.streamlit/secrets.toml` et y ajouter `GEMINI_API_KEY = 'votre_clé'` pour lancer l'analyse.")
        # Permet à l'utilisateur de voir l'interface même sans clé, mais l'analyse ne fonctionnera pas
        # st.stop() 


    # --- SECTION D'UPLOAD (Seule) ---
    uploaded_files = st.file_uploader("📁 **Chargement de CV :**", 
                                    type=["docx", "pdf", "pptx"], 
                                    accept_multiple_files=True,
                                    label_visibility="visible") 


    if uploaded_files:
        
        st.markdown("<h3>Résultats de l'Analyse</h3>", unsafe_allow_html=True)
        
        for uploaded in uploaded_files:
            
            st.markdown("<hr style='border: 1px solid #DDDDDD; margin: 1.5em 0;'>", unsafe_allow_html=True)
            
            col_file, col_btn_analyze, col_btn_docx, col_btn_pptx = st.columns([0.4, 0.2, 0.2, 0.2])
            
            with col_file:
                st.markdown(f"**📄 {uploaded.name}**")
                
            # --- Extraction du texte ---
            b = uploaded.read()
            text = ""
            extraction_error = None
            try:
                if uploaded.name.lower().endswith(".docx"):
                    text = extract_from_docx_bytes(b)
                elif uploaded.name.lower().endswith(".pdf"):
                    if not pdfplumber:
                        st.warning(f"Pour analyser {uploaded.name} (PDF), installez : `pip install pdfplumber`")
                        continue
                    text = extract_from_pdf_bytes(b)
                elif uploaded.name.lower().endswith(".pptx"):
                    if not Presentation:
                        st.warning(f"Pour analyser {uploaded.name} (PPTX), installez : `pip install python-pptx`")
                        continue
                    text = extract_from_pptx_bytes(b)
                else:
                    st.error(f"Fichier non pris en charge : {uploaded.name}")
                    continue
                
                st.session_state["extracted_texts"][uploaded.name] = text
                
            except Exception as e:
                extraction_error = f"Erreur d'extraction du texte pour {uploaded.name}: {e}"
                st.error(extraction_error)
                continue

            with col_btn_analyze:
                if st.button(f"🧠 Analyser", key=f"analyze_{uploaded.name}"):
                    if extraction_error:
                        st.error(f"Impossible de lancer l'analyse. Erreur d'extraction : {extraction_error}")
                    #elif not api_key:
                        #st.error("Impossible de lancer l'analyse. Configurez la clé API.")
                    else:
                        with st.spinner(f"Analyse de {uploaded.name} en cours par Gemini..."):
                            current_text = st.session_state["extracted_texts"].get(uploaded.name, text)
                            result = parse_with_llm(current_text)
                        
                        if "_error" in result:
                            st.error(result["_error"])
                        else:
                            # Mise à jour des résultats de session
                            new_results = [(fname, data) for fname, data in st.session_state["results"] if fname != uploaded.name]
                            new_results.append((uploaded.name, result))
                            st.session_state["results"] = new_results
                            st.success(f"✅ Analyse de {uploaded.name} réussie.")
                            st.rerun() # Rafraîchir pour afficher les boutons de téléchargement

            # --- Boutons de Téléchargement (Séparés) ---
            current_result = next(((fname, data) for fname, data in st.session_state["results"] if fname == uploaded.name), None)
            
            if current_result:
                fname, data = current_result
                
                # --- BOUTON DOCX (Word) ---
                with col_btn_docx:
                    buffer_docx = io.BytesIO()
                    output_name_docx = os.path.splitext(fname)[0] + "_GT.docx"
                    
                    try:
                        build_standard_docx(data, buffer_docx)
                        
                        st.download_button(
                            f"⬇️ Word (DOCX)",
                            data=buffer_docx.getvalue(),
                            file_name=output_name_docx,
                            mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
                            key=f"download_docx_{fname}"
                        )
                        
                    except Exception as e:
                        st.error(f"Erreur DOCX : {e}")

            
            else:
                # Assurer l'alignement lorsque les boutons de téléchargement n'apparaissent pas encore
                with col_btn_docx:
                    st.markdown("<div style='height: 38px;'></div>", unsafe_allow_html=True)
if __name__ == "__main__":
    pass  # laisse vide si tout ton code est déjà en haut
            