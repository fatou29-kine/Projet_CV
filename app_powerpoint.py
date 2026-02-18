import streamlit as st
import google.generativeai as genai
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
# Importation spécifique pour gérer l'indentation
from pptx.util import Inches 

from docx import Document
from pypdf import PdfReader
from io import BytesIO
import os
import re
import json
from pydantic import BaseModel, Field
from typing import List
from copy import deepcopy
import base64
def main():
    # --- CONFIG ---
    st.set_page_config(page_title="GTT CV Builder",page_icon="image.png", layout="wide")

    # Initialiser l'état de session si ce n'est pas déjà fait
    if 'cv_data' not in st.session_state:
        st.session_state.cv_data = None
    if 'cv_text' not in st.session_state:
        st.session_state.cv_text = None

    # ---------- DESIGN STREAMLIT (logo + CSS + header) ----------
    # NOTE: Le chemin vers le logo doit être valide pour que le code fonctionne localement.
    LOGO_PATH = "logo GTT nEw.png" 

    def get_base64_of_image(image_file):
        # Vérification d'existence du fichier
        if not os.path.exists(image_file):
            # Si le fichier n'existe pas, utiliser un placeholder pour ne pas bloquer l'application
            return "" 
        try:
            with open(image_file, "rb") as f:
                return base64.b64encode(f.read()).decode()
        except Exception:
            return ""

    _logo_b64 = get_base64_of_image(LOGO_PATH)

    PRIMARY = "#5C2D91"
    SECONDARY = "#4B9E4D"
    ACCENT = "#F3E6FF"
    BG = "#FAFAFA"

    st.markdown(f"""
    <style>
    /* page */
    .stApp {{ background-color: {BG}; font-family: 'Segoe UI', sans-serif; }}
    #MainMenu {{visibility: hidden;}} footer {{visibility: hidden;}} header {{visibility: hidden;}}

    /* header bande fixe */
    .header-band {{
        width: 100%;
        background-color: {PRIMARY};
        padding: 12px 24px;
        display: flex;
        align-items: center;
        gap: 16px;
        box-shadow: 0 2px 6px rgba(0,0,0,0.1);
        position: fixed;
        top: 0;
        left: 0;
        z-index: 9999;
    }}
    .header-band img {{
        height: 44px;
        width: auto;
        background-color: white; /* Fond clair pour le logo */
        padding: 4px; /* Un petit espace autour */
        border-radius: 6px; 
    }}
    .header-band h1 {{
        color: white;
        font-size: 1.5rem;
        margin: 0;
        font-weight: 800;
    }}

    /* Décaler le contenu sous le header */
    body .block-container {{ padding-top: 80px; }}

    /* uploader */
    div[data-testid="stFileUploader"] > div:first-child {{ font-size:1.05rem; font-weight:600; color:{PRIMARY}; }}
    div[data-testid="stFileUploader"] {{
        background: {ACCENT}; padding: 1.2rem; border-radius: 10px;
        border: 2px dashed {PRIMARY}55; margin-bottom: 1rem;
    }}

    /* buttons */
    .stButton>button {{ background-color: {PRIMARY}; color: white; border-radius: 8px; padding: 0.6rem 1.2rem; font-weight:700; }}
    .stButton>button:hover {{ background-color: #43206a; }}
    .stDownloadButton>button {{ background-color: {SECONDARY}; color: white; border-radius: 8px; padding: 0.6rem 1.2rem; font-weight:700; }}

    /* Style pour l'aperçu de l'analyse (comme sur l'image) */
    .data-box {{
        border: 1px solid #ddd;
        border-radius: 8px;
        padding: 15px;
        margin-bottom: 20px;
        background-color: white;
    }}
    .data-box h3 {{
        color: {PRIMARY};
        border-bottom: 2px solid {PRIMARY};
        padding-bottom: 5px;
        margin-top: 0;
        font-size: 1.2rem;
    }}
    .data-box ul {{
        list-style-type: disc;
        padding-left: 20px;
    }}
    </style>
    """, unsafe_allow_html=True)

    if _logo_b64:
        st.markdown(f"""
        <div class="header-band">
            <img src="data:image/png;base64,{_logo_b64}" alt="logo">
            <h1>Analyse & Standardisation CV</h1>
        </div>
        """, unsafe_allow_html=True)

    st.markdown("Génération du format Grant Thornton Technologies")
    st.divider()

    # ---------- COULEURS PPTX ----------
    MAUVE_FONCE = RGBColor(102, 0, 153)
    MAUVE_CLAIR = RGBColor(230, 214, 245)
    NOIR = RGBColor(0, 0, 0)
    FONT_NAME = "Trebuchet MS"

    # ---------- API GEMINI ----------
    try:
        # Utiliser st.secrets pour la clé API
        # NOTE : Pour l'exécution locale, vous devrez créer un fichier secrets.toml 
        # contenant GEMINI_API_KEY = "VOTRE_CLÉ_API"
        genai.configure(api_key=st.secrets["GEMINI_API_KEY"])
    except Exception:
        # Dans l'environnement de l'IA, cette partie est gérée différemment.
        # Pour le code local, s'assurer que la clé est bien configurée.
        pass 
        
    GEMINI_MODEL = "gemini-2.5-flash"

    # ---------- SCHEMA Pydantic ----------
    class Experience(BaseModel):
        company: str
        position: str
        period: str

    class Project(BaseModel):
        period: str
        organization: str
        country: str
        summary: str

    class CVData(BaseModel):
        NOM: str
        POSTE: str
        DOMAINE_D_EXPERTISE_SPECIFIQUE: str = Field(alias="DOMAINE D’EXPERTISE SPECIFIQUE")
        FORMATION: str
        PROFIL: str
        CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES: str
        
        REFERENCES_PERTINENTES: str = Field(
            description="Liste des références professionnelles ou des organisations/clients pertinents mentionnés dans le CV. Si cette section est absente ou vide, extraire la liste des clients ou des organisations mentionnés sur le CV (e.g. CNOSP, MEN, BAD, Commission UEMOA)."
        )
        
        experiences: List[Experience]
        projects: List[Project]

    # ---------- EXTRACTION CV ----------
    def extract_cv_text(file):
        text = ""
        try:
            if file.type == "application/pdf":
                reader = PdfReader(file)
                for page in reader.pages:
                    text += page.extract_text() or ""
                    text += "\n\n--- Nouvelle Page ---\n\n"
            elif file.type == "application/vnd.openxmlformats-officedocument.wordprocessingml.document":
                doc = Document(file)
                for para in doc.paragraphs:
                    text += para.text + "\n"
                for table in doc.tables:
                    text += "\n--- Contenu de Tableau (DOCX) ---\n"
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                        text += row_text + "\n"
                    text += "--- Fin de Tableau (DOCX) ---\n"
            elif file.type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
                st.warning("L'extraction de texte depuis un fichier PPTX est moins fiable que PDF/DOCX. Les données peuvent être incomplètes.")
                prs = Presentation(file)
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text_frame.text + "\n"
                        if shape.has_table:
                            for row in shape.table.rows:
                                row_text = " | ".join(cell.text.strip().replace("\n", " ") for cell in row.cells)
                                text += row_text + "\n"
                            text += row_text + "\n" # Pour s'assurer qu'un saut de ligne est après
                
            return re.sub(r'\s{3,}', '\n\n', text).strip()
        except Exception as e:
            st.error(f"Erreur lecture CV : {e}")
            return None

    # ---------- ANALYSE CV AVEC GEMINI (PROMPT MIS À JOUR : PLUS DE LIMITE À 4) ----------
    @st.cache_resource(show_spinner="Analyse du CV avec Gemini...")
    def analyze_cv_with_gemini(text):
        """Appelle Gemini pour analyser le texte et le valider contre le schéma Pydantic. Prompt plus précis pour PROFIL/FORMATION/DOMAINE."""
        prompt = f"""
    Tu es expert RH sénégalais. Extrais TOUT le contenu du CV suivant au format JSON strict. 

    Instructions cruciales pour l'extraction :
    1.  **PROFIL**: Extrais le profil professionnel. Limitation stricte à 4 phrases maximum. Sépare chaque phrase par un caractère de **retour à la ligne unique ('\n')** pour forcer l'affichage vertical (une phrase par ligne). Rédige un profil concis et percutant.
    2.  **DOMAINE D’EXPERTISE SPECIFIQUE**: Extrais **TOUS** les domaines d'expertise pertinents listés dans le CV. Sépare chaque domaine par un caractère de **retour à la ligne unique ('\n')** dans le JSON.
    3.  **FORMATION**: Extrais **TOUTES** les formations ou diplômes pertinents listés dans le CV. Sépare chaque formation par un caractère de **retour à la ligne unique ('\n')** dans le JSON.
    4. **REFERENCES_PERTINENTES**: Extrais **TOUTES** les références professionnelles, clients ou organisations mentionnés dans le CV. Sépare chaque référence par un caractère de **retour à la ligne unique ('\n')** dans le JSON.
    5. **CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES**: Extrais **TOUTES** les certifications professionnelles pertinentes listées dans le CV. Sépare chaque certification par un caractère de **retour à la ligne unique ('\n')** dans le JSON.


    Extrais les champs suivants :
    - NOM, POSTE, PROFIL, DOMAINE D’EXPERTISE SPECIFIQUE
    - FORMATION, CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES, REFERENCES_PERTINENTES
    Pour le champ REFERENCES_PERTINENTES : 
    - Extrais toutes les références professionnelles ou organisations mentionnées dans le CV. 
    - Si une mission ou un rôle associé est mentionné, ajoute-la **après un double point**. 
    - Chaque ligne doit respecter le format exact : "Référence : Mission(Année ou Période)". 
    - Extrais la période (Année, Années, ou Période Complète) si elle est associée à la référence et ajoute-la entre parenthèses à la fin.
    - Si aucune mission n’est mentionnée pour une référence, mets seulement le nom de la référence.
    - experiences : toutes les expériences (company, position, period)
    - projects : tous les projets (period, organization, country, summary)


    CV :

    {text[:49000]}
    """
        try:
            model = genai.GenerativeModel(GEMINI_MODEL,
                generation_config=genai.GenerationConfig(
                    response_mime_type="application/json",
                    response_schema=CVData,
                    temperature=0.0
                ))
            response = model.generate_content(prompt)
            raw = response.text.strip()
            
            # Nettoyage pour les blocs de code Markdown
            if raw.startswith("```"): 
                raw = raw.split("```",1)[1].rsplit("```",1)[0]
            if raw.startswith("json\n"):
                raw = raw[5:]

            return CVData.model_validate(json.loads(raw))
        except json.JSONDecodeError as e:
            st.error(f"Erreur de décodage JSON de Gemini : Le modèle n'a pas renvoyé un JSON valide. {e}")
            st.code(raw)
            return None
        except Exception as e:
            st.error(f"Erreur d'API ou de validation Pydantic : {e}")
            return None

    # ---------- FONCTIONS PPTX (Fonctions de support pour la génération PPTX) ----------

    # --- INSERTION LOGO PPTX (MODIFIÉE POUR 6 LOGOS) ---
    def insert_logos_on_first_slide(prs, logo_streams):
        """Insère une liste de logos (streams d'octets) sur la première diapositive."""
        if not prs.slides:
            st.warning("Aucune diapositive trouvée pour insérer le logo.")
            return
            
        first_slide = prs.slides[0]
        
        # Paramètres de positionnement et taille pour 6 logos
        START_LEFT = Inches(0.4)
        TOP = Inches(4.7) 
        LOGO_WIDTH = Inches(0.8) # Taille réduite pour en faire tenir 6
        SPACING = Inches(0.1)
        
        current_left = START_LEFT
        
        for i, logo_stream in enumerate(logo_streams):
            try:
                # Ajoute l'image à partir du stream d'octets (BytesIO)
                first_slide.shapes.add_picture(logo_stream, current_left, TOP, width=LOGO_WIDTH)
                
                # Calcule la position pour le logo suivant
                current_left += LOGO_WIDTH + SPACING
                
                # Limite stricte à 6 logos (indices 0 à 5)
                if i >= 5: 
                    break
                    
            except Exception as e:
                st.warning(f"Erreur lors de l'insertion du logo {i+1}: {e}")
                
        st.success(f"Tentative d'insertion de {min(len(logo_streams), 6)} logos sur la première diapositive.")
    # --- FIN INSERTION LOGO ---
    # --- FIN INSERTION LOGO ---

    def get_template_info(template_path):
        prs = Presentation(template_path)
        placeholders = {}
        exp_slide = None
        exp_table = None
        proj_slides = []

        placeholder_keys = {
            "NOM": ["NOM"],
            "POSTE": ["POSTE"],
            "DOMAINE D’EXPERTISE SPECIFIQUE": ["DOMAINE D’EXPERTISE SPECIFIQUE", "DOMAINE D'EXPERTISE SPECIFIQUE", "DOMAINE D EXPERTISE SPECIFIQUE"],
            "FORMATION": ["FORMATION"],
            "PROFIL": ["PROFIL"],
            "CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES": ["CERTIFICATIONS PROFESSIONNELLES PERTINENTES", "CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES", "CERTIFICATIONS"],
            "REFERENCES_PERTINENTES": ["REFERENCES PERTINENTES", "REFERENCES_PERTINENTES", "REFERENCES"]
        }

        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    txt = shape.text_frame.text.upper().replace("_", " ")
                    for key, variants in placeholder_keys.items():
                        if any(var.upper() in txt for var in variants):
                            placeholders.setdefault(key, []).append(shape)
                if shape.has_table:
                    if not shape.table.rows or not shape.table.rows[0].cells:
                        continue
                    header = " ".join(c.text.lower() for c in shape.table.rows[0].cells)
                    if any(k in header for k in ["société", "company", "poste", "période", "emploi"]):
                        if not exp_slide:
                            exp_slide = slide
                            exp_table = shape.table
                    if any(k in header for k in ["organisation", "country", "pays", "résumé", "summary"]):
                        proj_slides.append(slide)

        return placeholders, exp_slide, exp_table, proj_slides, prs

    def insert_row(table):
        if not table.rows:
            raise IndexError("Tableau corrompu : Impossible de trouver une ligne modèle.")
        new_row_index = len(table.rows)
        # Copie la structure de la première ligne comme modèle
        template_row_tr = table.rows[0]._tr
        tr = deepcopy(template_row_tr)
        # Insère la nouvelle ligne
        table._tbl.append(tr)
        # Force la mise à jour de la liste des lignes (important après ajout)
        table._rows = None
        return table.rows[new_row_index]

    def duplicate_slide(prs, slide):
        # Créer une nouvelle diapositive avec le même layout
        new_slide = prs.slides.add_slide(slide.slide_layout)
        # Copier les formes (shapes) de la diapositive source vers la nouvelle
        for shape in slide.shapes:
            el = deepcopy(shape.element)
            new_slide.shapes._spTree.insert_element_before(el, "p:extLst")
        return new_slide

    def fill_experiences(prs, template_slide, template_table, experiences):
        if not template_slide or not experiences: return
        # Trouver l'index de la diapositive d'expérience pour savoir où insérer les duplications
        slide_id = next((i for i, slide in enumerate(prs.slides) if slide == template_slide), None)
        if slide_id is None: return
        
        idx = 0
        while idx < len(experiences):
            # Utiliser la diapo template pour la première itération, puis dupliquer pour les suivantes
            slide = prs.slides[slide_id] if idx == 0 else duplicate_slide(prs, template_slide)
            
            # Trouver le tableau sur la diapositive
            table = next((s.table for s in slide.shapes if s.has_table), None)
            if not table: 
                idx += 10 # Passer à la diapo suivante (si on a sauté la duplication)
                continue
                
            # Nettoyer les lignes existantes du tableau (sauf l'en-tête)
            table._rows = None
            while len(table.rows) > 1:
                table._tbl.remove(table.rows[1]._tr)
                table._rows = None
                
            # Remplir jusqu'à 10 lignes par diapositive
            for _ in range(10):
                if idx >= len(experiences): break
                row = insert_row(table)
                row.height = Inches(0.45)   # AUGMENTE LA HAUTEUR ICI (0.45 = encore plus grand)
                if len(row.cells) >= 3:
                    # Colonne 1: Société/Company
                    row.cells[0].text = experiences[idx].company or ""
                    # Colonne 2: Poste/Position
                    row.cells[1].text = experiences[idx].position or ""
                    # Colonne 3: Période
                    row.cells[2].text = experiences[idx].period or ""
                    
                    for cell in row.cells:
                        # Appliquer le fond clair (Mauve clair)
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = MAUVE_CLAIR
                        
                        for paragraph in cell.text_frame.paragraphs:
                            # Rétablir le niveau de retrait à 0 (pour les tableaux)
                            paragraph.level = 0
                            
                            # Réinitialiser les indentations pour éviter les décalages
                            if hasattr(paragraph, 'paragraph_format'):
                                paragraph.paragraph_format.left_indent = Inches(0)
                                paragraph.paragraph_format.first_line_indent = Inches(0)
                                
                            # S'assurer qu'un run existe pour le formatage si le texte est vide
                            if not paragraph.runs:
                                paragraph.text = cell.text
                            
                            # Appliquer le formatage de police
                            for run in paragraph.runs:
                                run.font.name = FONT_NAME
                                run.font.size = Pt(11) # Taille standard (11pt)
                                run.font.color.rgb = NOIR
                                run.font.bold = False
                idx += 1
    def fill_projects(prs, slides, projects):
        i = 0
        for slide in slides:

            # Chercher le tableau dans la slide
            table = next((s.table for s in slide.shapes if s.has_table), None)
            if not table or i >= len(projects):
                continue

            # Réinitialisation des lignes existantes
            table._rows = None
            while len(table.rows) > 1:
                table._tbl.remove(table.rows[1]._tr)
                table._rows = None

            # Insérer jusqu'à 3 projets par slide
            for _ in range(3):
                if i >= len(projects):
                    break

                row = insert_row(table)
                row.height = Inches(0.60)

                if row and len(row.cells) >= 4:

                    # Remplissage des colonnes classiques
                    row.cells[0].text = projects[i].period
                    row.cells[1].text = projects[i].organization
                    row.cells[2].text = projects[i].country

                    # ============================
                    #  SUMMARY : TOUT EN PUCES
                    # ============================
                    summary_cell = row.cells[3]
                    summary_cell.text = ""  # vider avant de reconstruire

                    raw_summary = projects[i].summary.strip()
                    # Découper sur les retours à la ligne ou les points suivis d'un espace
                    lines = re.split(r'\.\s+', raw_summary)
                    lines = [l.strip() for l in lines if l.strip()]
                    tf = summary_cell.text_frame
                    tf.clear()

                    for l in lines:
                        p = tf.add_paragraph()
                        p.level = 0
                        p.text = "• " + l
                        if p.runs:
                            for r in p.runs:
                                r.font.size = Pt(11)
                                r.font.bold = False
                                r.font.name = 'Trebuchet MS'
                                r.font.color.rgb = NOIR

                    # ============================
                    # STYLE GENERAL DES CELLULES (toutes)
                    # ============================
                    for cell in row.cells:
                        cell.fill.solid()
                        cell.fill.fore_color.rgb = MAUVE_CLAIR
                        for p in cell.text_frame.paragraphs:
                            p.level = 0
                            if hasattr(p, 'paragraph_format'):
                                p.paragraph_format.left_indent = Inches(0)
                                p.paragraph_format.first_line_indent = Inches(0)
                            for r in p.runs:
                                r.font.size = Pt(11)
                                r.font.name = 'Trebuchet MS'
                                r.font.bold = False  # DESACTIVER GRAS
                                r.font.color.rgb = NOIR

                i += 1
    def apply_formatting(prs):
        """
        Applique le formatage final uniquement aux titres statiques spécifiques 
        qui ne sont pas remplis par l'IA (pour ne pas écraser le formatage inséré).
        """
        STATIC_SPECIAL_TEXT = ["CV DÉTAILLÉ", "DIRECTEUR DE MISSION", "ÉQUIPE D’INTERVENTION"] 
        
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame: continue
                
                text = " ".join(p.text for p in shape.text_frame.paragraphs).upper()
                
                # Formatage SPÉCIAL pour les TITRES STATIQUES (ceux qui ne sont pas remplis par le CV)
                if any(x in text for x in STATIC_SPECIAL_TEXT): 
                    for p in shape.text_frame.paragraphs:
                        p.alignment = PP_ALIGN.CENTER
                        for run in p.runs:
                            run.font.name = FONT_NAME
                            run.font.size = Pt(18)
                            run.font.color.rgb = MAUVE_FONCE
                    continue 


    # --------------------------------------------------------------------------------------
    # ---------- LOGIQUE DE SÉLECTION PAR CHECKBOXES ---------------------------------------
    # --------------------------------------------------------------------------------------

    def display_checkbox_selection(title, data_string, session_key, columns=1):
        """Affiche les éléments extraits sous forme de cases à cocher et gère leur état. Initialisation à False (décoché)."""
        # Sépare la chaîne d'éléments par saut de ligne ('\n')
        items = [item.strip() for item in data_string.split('\n') if item.strip()]
        st.markdown(f"### {title}")
        
        # 1. Initialiser l'état de session : tout est DÉCOCHÉ (False) par défaut.
        # On vérifie si l'état actuel correspond aux éléments extraits du CV, sinon on réinitialise
        current_items_set = set(items)
        current_selection_keys = set(st.session_state.get(session_key, {}).keys())

        # Réinitialise si la clé n'existe pas ou si les éléments du CV ont changé (nouveau CV)
        if session_key not in st.session_state or not isinstance(st.session_state[session_key], dict) or current_items_set != current_selection_keys:
            # Initialisation par défaut à False pour tous les nouveaux éléments
            st.session_state[session_key] = {item: False for item in items}
        
        current_selection = st.session_state[session_key]
        new_selection = {}
        
        # Utilisation de colonnes Streamlit pour un affichage plus propre
        cols = st.columns(columns)
        
        for i, item in enumerate(items):
            col = cols[i % columns]
            
            # Récupère l'état précédent ou FALSE par défaut.
            is_checked = current_selection.get(item, False) 
            
            with col:
                # Créer la case à cocher
                # Utilisation d'un key unique et stable pour chaque élément
                key_val = f"{session_key}_{i}_{hash(item)}" 
                new_selection[item] = st.checkbox(item, 
                                                value=is_checked, 
                                                key=key_val, 
                                                help=item)

        # 2. Mettre à jour l'état de session avec les nouvelles sélections
        st.session_state[session_key] = new_selection
        
        # 3. Retourner la liste des éléments sélectionnés
        selected_items = [item for item, selected in st.session_state[session_key].items() if selected]
        return selected_items


    def get_selected_content(session_key, fallback_data):
        """Fonction utilitaire pour récupérer le contenu sélectionné, ou le contenu brut si non sélectionné/analysé."""
        selected_dict = st.session_state.get(session_key)
        
        # Si la sélection existe et n'est pas vide (au moins un élément coché)
        # On vérifie si l'état contient des clés, sinon on retourne vide
        if selected_dict is not None and isinstance(selected_dict, dict) and any(selected_dict.values()):
            # Joindre uniquement les éléments marqués 'True' avec un retour à la ligne
            selected_items = [item for item, selected in selected_dict.items() if selected]
            return '\n'.join(selected_items)
        
        # Si la sélection n'existe pas ou si rien n'a été coché, on retourne une chaîne vide.
        return "" 

    # --------------------------------------------------------------------------------------
    # ---------- LOGIQUE D'AFFICHAGE D'ANALYSE (MISE À JOUR) --------------------------------
    # --------------------------------------------------------------------------------------

    def display_analysis_preview(data: CVData):
        st.markdown("## 🔍 Aperçu de l'Analyse du CV et Sélection pour le PPTX")
        st.markdown("⚠️ **Veuillez cocher les éléments que vous souhaitez inclure dans le PPTX.** (Ils sont tous décochés par défaut, sauf le Profil qui est obligatoire)")
        st.markdown(f"**Nom:** {data.NOM} | **Poste:** {data.POSTE}")
        st.markdown("---")
        
        # Première Ligne : Profil / Références Pertinentes
        col_l1, col_r1 = st.columns(2)

        with col_l1:
            st.markdown('<div class="data-box">', unsafe_allow_html=True)
            st.markdown("### Profil") 
            profile_content = st.text_area("Modifier le profil :", data.PROFIL, height=150)
            st.session_state.cv_data.PROFIL = profile_content
            st.markdown('</div>', unsafe_allow_html=True)

        with col_r1:
            st.markdown('<div class="data-box">', unsafe_allow_html=True)
            # Références Pertinentes (Case à cocher)
            display_checkbox_selection("Références Pertinentes / Clients", 
                                    data.REFERENCES_PERTINENTES, 
                                    'selected_references',
                                    columns=1)
            st.markdown('</div>', unsafe_allow_html=True)

        st.markdown("---") 
        
        # Deuxième Ligne : Domaine d'Expertise / Formation & Certifications
        col_l2, col_r2 = st.columns(2)
        
        with col_l2:
            st.markdown('<div class="data-box">', unsafe_allow_html=True)
            # Domaines d'Expertise (Case à cocher)
            display_checkbox_selection("Domaine d'Expertise Spécifique", 
                                    data.DOMAINE_D_EXPERTISE_SPECIFIQUE, 
                                    'selected_domaines',
                                    columns=1)
            st.markdown('</div>', unsafe_allow_html=True)

        with col_r2:
            st.markdown('<div class="data-box">', unsafe_allow_html=True)
            # Formation (Case à cocher)
            display_checkbox_selection("Formation", 
                                    data.FORMATION, 
                                    'selected_formations',
                                    columns=1)
            
            st.markdown("<hr style='border-top: 1px dashed #ccc;'>", unsafe_allow_html=True)
            
            # Certifications (Case à cocher)
            display_checkbox_selection("Certifications Professionnelles Pertinentes", 
                                    data.CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES, 
                                    'selected_certifications',
                                    columns=1)
            st.markdown('</div>', unsafe_allow_html=True)
        
        st.markdown("---")
        
        # Affichage des expériences professionnelles
        st.markdown("### Expériences Professionnelles (Extrait)")
        exp_data = [{'Période': e.period, 'Société': e.company, 'Poste': e.position} for e in data.experiences]
        st.dataframe(exp_data, use_container_width=True, hide_index=True)
        
        # Affichage des projets
        st.markdown("### Projets Réalisés (Extrait)")
        proj_data = [{'Période': p.period, 'Organisation': p.organization, 'Pays': p.country, 'Résumé': p.summary[:100] + '...'} for p in data.projects]
        st.dataframe(proj_data, use_container_width=True, hide_index=True)
        
    # ---------- MAIN LOGIC (CORRIGÉE) ----------
    template_path = "CV PPT template.pptx" 
    if not os.path.exists(template_path):
        st.error(f"Template PPTX manquant : {template_path} ! Assurez-vous d'avoir le fichier 'CV PPT template.pptx' dans le même répertoire.")

    # --- FIN UPLOAD LOGO ---

    cv_file = st.file_uploader("📁Téléchargez votre CV", type=["pdf","docx","pptx"], key="cv_uploader")

    col1, col2 = st.columns(2)
    # --- UPLOAD LOGO ---
    logo_files = st.file_uploader("🖼️ Téléchargez les logos à insérer", type=["png","jpg","jpeg"],accept_multiple_files=True, key="logo_uploader")


    # --- BOUTON D'ANALYSE ---
    if cv_file:
        if st.button("🔍Afficher l'Analyse CV"):
            st.session_state.cv_text = extract_cv_text(cv_file)
            if st.session_state.cv_text:
                # Vider l'état des sélections précédentes pour forcer la réinitialisation
                st.session_state['selected_references'] = {}
                st.session_state['selected_domaines'] = {}
                st.session_state['selected_formations'] = {}
                st.session_state['selected_certifications'] = {}
                
                # Re-exécution de l'analyse 
                st.session_state.cv_data = analyze_cv_with_gemini(st.session_state.cv_text) 
                
                if st.session_state.cv_data:
                    st.success("Analyse réussie ! Visualisez l'aperçu ci-dessous et effectuez vos sélections.")
                    st.balloons()
                else:
                    st.error("L'analyse du CV par Gemini a échoué.")
            else:
                st.error("Impossible d'extraire le texte du CV.")
                
        # ✅ L'AJOUT CRUCIAL POUR AFFICHER LES CHECKBOXES ET MAINTENIR LEUR ÉTAT
        if st.session_state.cv_data:
            display_analysis_preview(st.session_state.cv_data)
            st.markdown("---")  # ligne de séparation pour faire apparaître le bouton PPTX en bas   
        # --- BOUTON DE GÉNÉRATION PPTX ---
        if st.button("⚡Générer le PPTX"):
            # Tenter d'analyser si cela n'a pas été fait (ou si la session a expiré)
            if not st.session_state.cv_data and cv_file:
                st.warning("Analyse en cours... Veuillez patienter.")
                st.session_state.cv_text = extract_cv_text(cv_file)
                if st.session_state.cv_text:
                    st.session_state.cv_data = analyze_cv_with_gemini(st.session_state.cv_text)
            
            data = st.session_state.cv_data
            
            if data:
                
                # --- RÉCUPÉRATION DES DONNÉES FILTRÉES PAR L'UTILISATEUR ---
                # Si aucune case n'est cochée, la fonction get_selected_content renvoie une chaîne vide ""
                filtered_domaines = get_selected_content('selected_domaines', data.DOMAINE_D_EXPERTISE_SPECIFIQUE)
                filtered_references = get_selected_content('selected_references', data.REFERENCES_PERTINENTES)
                filtered_formations = get_selected_content('selected_formations', data.FORMATION)
                filtered_certifications = get_selected_content('selected_certifications', data.CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES)

                # --- FIN LOGIQUE DE SÉLECTION ---
                
                st.success("✅CV analysé ! Démarrage de la génération PPTX...")
                
                try:
                    placeholders, exp_slide, exp_table, proj_slides, prs = get_template_info(template_path)

                    # --- APPEL INSERTION LOGO ---
                    if logo_files :
                        # Utiliser BytesIO pour passer le contenu de l'image à la fonction sans enregistrement local
                        logo_streams = [BytesIO(f.read()) for f in logo_files]
                        insert_logos_on_first_slide(prs, logo_streams)
                    # --- FIN APPEL INSERTION LOGO ---

                    # MISE À JOUR DU MAPPING AVEC LES DONNÉES FILTRÉES
                    mapping = {
                        "NOM": data.NOM,
                        "POSTE": data.POSTE,
                        "DOMAINE D’EXPERTISE SPECIFIQUE": filtered_domaines,
                        "FORMATION": filtered_formations,
                        "PROFIL": data.PROFIL, # Le profil est toujours inclus
                        "CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES": filtered_certifications,
                        "REFERENCES_PERTINENTES": filtered_references
                    }
                    
                    pydantic_to_pptx_token = {
                        "NOM": "NOM",
                        "POSTE": "POSTE",
                        "DOMAINE D’EXPERTISE SPECIFIQUE": "DOMAINE D’EXPERTISE SPECIFIQUE",
                        "FORMATION": "FORMATION",
                        "PROFIL": "PROFIL",
                        "CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES": "CERTIFICATIONS PROFESSIONNELLES PERTINENTES",
                        "REFERENCES_PERTINENTES": "REFERENCES PERTINENTES"
                    }

                    MAIN_TITLES_KEYS = ["NOM", "POSTE"]
                    SECTION_HEADERS_KEYS = [
                        "PROFIL", 
                        "DOMAINE D’EXPERTISE SPECIFIQUE", 
                        "REFERENCES_PERTINENTES", 
                        "CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES",
                        "FORMATION" 
                    ]
                    
                    DISPLAY_TITLES = {
                        "PROFIL": "PROFIL",
                        "DOMAINE D’EXPERTISE SPECIFIQUE": "DOMAINE D'EXPERTISE SPÉCIFIQUE",
                        "REFERENCES_PERTINENTES": "RÉFÉRENCES PERTINENTES",
                        "CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES": "CERTIFICATIONS PROFESSIONNELLES PERTINENTES",
                        "FORMATION": "FORMATION",
                    }


                    with st.spinner("Remplissage du template PPTX..."):
                        for ph_key, shapes in placeholders.items():
                            value = mapping.get(ph_key, "").strip()
                            token = pydantic_to_pptx_token.get(ph_key, ph_key)
                            
                            for s in shapes:
                                if s.has_text_frame and token in s.text:
                                    
                                    # 1. Gestion des Titres Principaux (NOM, POSTE)
                                    if ph_key in MAIN_TITLES_KEYS:
                                        # Remplacer le token directement
                                        s.text = s.text.replace(token, str(value))
                                        
                                        # DÉTERMINATION DE LA TAILLE DE POLICE ADAPTATIVE
                                        font_size = Pt(12) # Taille par défaut pour le NOM (12pt)
                                        if ph_key == "POSTE":
                                            # Ajustement très léger pour le POSTE si très long
                                            if len(str(value)) > 35:
                                                font_size = Pt(12)
                                            else:
                                                font_size = Pt(12)
                                        
                                        # Appliquer le formatage aux paragraphes existants
                                        for p in s.text_frame.paragraphs:
                                            p.alignment = PP_ALIGN.CENTER
                                            if hasattr(p, 'paragraph_format'):
                                                p.paragraph_format.left_indent = Inches(0)
                                                p.paragraph_format.first_line_indent = Inches(0)
                                                
                                            for run in p.runs:
                                                run.font.name = FONT_NAME
                                                run.font.size = font_size
                                                run.font.color.rgb = NOIR
                                                run.font.bold = True
                                    
                                    # 2. Gestion des En-têtes de Section (PROFIL, DOMAINE, FORMATION, etc.) 
                                    elif ph_key in SECTION_HEADERS_KEYS:
                                        title_text = DISPLAY_TITLES.get(ph_key, ph_key).upper()
                                        tf = s.text_frame
                                        
                                        # a) Vider la zone de texte du placeholder (sauf le premier paragraphe qui devient le titre)
                                        while len(tf.paragraphs) > 1:
                                            p = tf.paragraphs[1]
                                            tf.paragraphs._element.remove(p._element)
                                        tf.paragraphs[0].text = "" # Vider le texte de l'ancien placeholder
                                        
                                        # b) Recréer le titre de section dans le premier paragraphe
                                        title_p = tf.paragraphs[0]
                                        title_p.text = title_text
                                        title_p.level = 0
                                        title_p.alignment = PP_ALIGN.LEFT
                                        
                                        # Formatage du titre
                                        for run in title_p.runs:
                                            run.font.name = FONT_NAME
                                            run.font.size = Pt(14) # Taille pour les titres de section
                                            run.font.color.rgb = NOIR
                                            run.font.bold = True
                                        
                                        # c) Insérer le contenu s'il n'est pas vide avec le formatage spécifique
                                        if value:
                                            items = [item.strip() for item in str(value).split('\n') if item.strip()]
                                            for line in items:
                                                p = tf.add_paragraph() # Ajoute un nouveau paragraphe pour chaque ligne de contenu
                                                
                                                # --- LOGIQUE DE FORMATAGE DES LIGNES (CORRIGÉE) ---
                                                
                                                # Cas 1 : PROFIL (Texte simple, aligné à gauche, pas de puce)
                                                if ph_key == "PROFIL":
                                                    p.text = line
                                                    p.level = 0 
                                                    p.alignment = PP_ALIGN.LEFT
                                                    
                                                # Cas 2 : CERTIFICATIONS (Puce '✔' insérée manuellement)
                                                elif ph_key == "CERTIFICATIONS_PROFESSIONNELLES_PERTINENTES":
                                                    p.text = "✔ " + line
                                                    p.level = 0 
                                                    p.alignment = PP_ALIGN.LEFT
                                                    
                                                # Cas 3 : Autres (Domaine, Formation, Références) (Puce '•' insérée manuellement)
                                                else:
                                                    p.text = "• " + line
                                                    p.level = 0 
                                                    p.alignment = PP_ALIGN.LEFT
                                                
                                                # --- Application du formatage d'indentation (important) ---
                                                # Réinitialiser l'indentation par défaut du niveau 0
                                                if hasattr(p, 'paragraph_format'): 
                                                    p.paragraph_format.left_indent = Inches(0)
                                                    p.paragraph_format.first_line_indent = Inches(0)
                                                
                                                # --- Application du formatage de police pour la ligne ---
                                                # S'assurer qu'un run existe pour appliquer le formatage
                                                if not p.runs:
                                                    # Ajoute le texte au paragraphe si pas de run (cela crée un run)
                                                    p.text = line if ph_key == "PROFIL" else p.text
                                                
                                                if p.runs:
                                                    run = p.runs[0]
                                                    run.font.name = FONT_NAME
                                                    run.font.size = Pt(11) # Taille standard pour le contenu (11pt)
                                                    run.font.color.rgb = NOIR
                                                    run.font.bold = False
                                                # --- FIN LOGIQUE DE FORMATAGE DES LIGNES ---
                                                
                    # --- Remplissage des Expériences et Projets ---
                    fill_experiences(prs, exp_slide, exp_table, data.experiences)
                    fill_projects(prs, proj_slides, data.projects)
                    
                    # --- Formatage final des titres statiques ---
                    apply_formatting(prs)

                    # --- Sauvegarde du fichier ---
                    output = BytesIO()
                    prs.save(output)
                    output.seek(0)
                    
                    st.download_button(
                        label="⬇️ Télécharger le CV PPTX",
                        data=output,
                        file_name=f"CV_{data.NOM.replace(' ', '_')}_GTT.pptx",
                        mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                        key="download_pptx_button"
                    )
                    st.balloons()

                except Exception as e:
                    st.error(f"Une erreur critique est survenue lors de la génération du PPTX : {e}")
                    st.exception(e)
if __name__ == "__main__":
    pass  # laisse vide si tout ton code est déjà en haut
